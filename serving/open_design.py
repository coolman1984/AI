"""Output / design layer (MASTER_PLAN K.1 layer 6).

Turns an **audited, signed-off** card into a visual artifact. Real today:
  - a professional one-A4 HTML dashboard with an inline SVG variance-bridge chart
  - a real PPTX management slide (python-pptx)
**Open Design** (nexu-io) plugs in via `OpenDesignRenderer` for richer/branded output
(dashboards/PDF/PPTX/MP4). Output ONLY — it renders numbers, never computes or invents them,
and only after sign-off.
"""
from __future__ import annotations

import importlib.util
from pathlib import Path


def _svg_bridge(parts, width: int = 460, row_h: int = 26) -> str:
    maxv = max((abs(p.variance) for p in parts), default=1) or 1
    mid = width / 2
    rows = []
    for i, p in enumerate(parts):
        y = i * row_h
        w = (abs(p.variance) / maxv) * (width / 2 - 70)
        x, color = (mid, "#d9534f") if p.variance >= 0 else (mid - w, "#5cb85c")
        rows.append(
            f'<rect x="{x:.0f}" y="{y + 4}" width="{w:.0f}" height="16" rx="2" fill="{color}"/>'
            f'<text x="6" y="{y + 17}" font-size="12" fill="#333">{p.dim_value}</text>'
            f'<text x="{width - 4}" y="{y + 17}" font-size="11" text-anchor="end" fill="#333">'
            f'{p.variance:+,.0f}</text>'
        )
    h = len(parts) * row_h + 8
    return (
        f'<svg viewBox="0 0 {width} {h}" width="100%" height="{h}" role="img">'
        f'<line x1="{mid}" y1="0" x2="{mid}" y2="{h}" stroke="#ccc"/>' + "".join(rows) + "</svg>"
    )


def render_dashboard_html(ctx: dict) -> str:
    card, bridge, audit = ctx["card"], ctx["bridge"], ctx["audit"]
    decomposition = ctx.get("decomposition")
    signoff = ctx.get("signoff")
    conf = card.confidence["confidence"]
    conf_color = {"High": "#5cb85c", "Medium": "#f0ad4e", "Low": "#d9534f"}.get(conf, "#777")
    nums = "".join(
        f'<div class="kpi"><div class="v">{n.value:+,.0f}</div>'
        f'<div class="l">{n.label}</div>'
        f'<div class="e">{n.evidence.source} · {n.evidence.method}</div></div>'
        for n in card.key_numbers
    )
    cites = "".join(f"<li>{e.source} — {e.locator} ({e.method})</li>" for e in card.evidence)
    questions = "".join(
        f"<div class='q'>❓ {q['question']}<ul>"
        + "".join(f"<li>{o}</li>" for o in q["options"]) + "</ul></div>"
        for q in audit.questions
    )
    sign = (
        f"<span class='ok'>✔ signed off by {signoff.approver}</span>"
        if (signoff and signoff.approved) else "<span class='no'>⛔ not signed — not released</span>"
    )
    why_now = ""
    if decomposition is not None and card.driver_split is not None:
        why_now = (
            "<div class=\"panel\"><b>Operating driver view (vs standard cost)</b>"
            f"<div class=\"kpis\">"
            f"<div class=\"kpi\"><div class=\"v\">{card.driver_split.price.value:+,.0f}</div><div class=\"l\">Price effect</div></div>"
            f"<div class=\"kpi\"><div class=\"v\">{card.driver_split.volume.value:+,.0f}</div><div class=\"l\">Volume effect</div></div>"
            f"<div class=\"kpi\"><div class=\"v\">{card.driver_split.mix.value:+,.0f}</div><div class=\"l\">Mix effect</div></div>"
            "</div>"
            f"<small>Standard-cost total {card.driver_split.total.value:+,.0f} · reconciles={decomposition.reconciles()}</small>"
            "</div>"
        )
    return f"""<!doctype html><meta charset=\"utf-8\">
<title>{card.headline}</title>
<style>
 body{{font-family:-apple-system,Segoe UI,Roboto,sans-serif;max-width:780px;margin:24px auto;color:#222}}
 .head{{border-left:6px solid {conf_color};padding:6px 14px;margin-bottom:14px}}
 h1{{font-size:20px;margin:0 0 4px}} .sub{{color:#666;font-size:13px}}
 .badge{{display:inline-block;background:{conf_color};color:#fff;border-radius:10px;padding:2px 10px;font-size:12px}}
 .kpis{{display:flex;gap:12px;margin:14px 0}}
 .kpi{{flex:1;border:1px solid #eee;border-radius:8px;padding:10px}}
 .kpi .v{{font-size:22px;font-weight:700}} .kpi .l{{font-size:12px;color:#444}} .kpi .e{{font-size:10px;color:#999;margin-top:4px}}
 .panel{{border:1px solid #eee;border-radius:8px;padding:12px;margin:12px 0}}
 .q{{background:#fff7e6;border:1px solid #f0ad4e;border-radius:6px;padding:8px;margin:6px 0;font-size:13px}}
 .ok{{color:#3c763d;font-weight:700}} .no{{color:#a94442;font-weight:700}}
 small,li{{font-size:12px;color:#555}}
</style>
<div class="head">
  <h1>{card.headline}</h1>
  <div class="sub">Decision: <b>{card.decision_needed}</b> · Period {ctx.get('period','')} ·
    Confidence <span class="badge">{conf}</span> · {sign}</div>
</div>
<div class="kpis">{nums}</div>
{why_now}
<div class="panel"><b>Variance bridge (by sub-assembly)</b>{_svg_bridge(bridge.parts)}</div>
<div class="panel"><b>Drivers:</b> {'; '.join(card.drivers)}<br>
  <b>Actions:</b> {'; '.join(card.actions)}
  {('<br><b>Risks:</b> ' + '; '.join(card.risks)) if card.risks else ''}</div>
<div class="panel"><b>Audit (independent re-run):</b> recomputed variance
  {audit.recomputed['variance']:+,.0f}, certainty {audit.certainty['overall']:.0%}
  {questions}</div>
<div class="panel"><b>Evidence</b><ul>{cites}</ul></div>
<small>Generated by Factory Second Brain · numbers from the data engine · audited &amp; human-signed.</small>
"""


class PptxRenderer:
    """Real PowerPoint export (python-pptx) — a management slide."""
    def available(self) -> bool:
        return importlib.util.find_spec("pptx") is not None

    def render(self, ctx: dict, out_path: str) -> str:
        from pptx import Presentation
        from pptx.util import Inches, Pt

        card = ctx["card"]
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        slide.shapes.title.text = card.headline
        box = slide.shapes.add_textbox(Inches(0.5), Inches(1.6), Inches(9), Inches(5))
        tf = box.text_frame
        tf.word_wrap = True
        tf.text = f"Decision: {card.decision_needed}  |  Confidence: {card.confidence['confidence']}"

        def add(text: str, size: int) -> None:
            para = tf.add_paragraph()
            para.text = text
            para.font.size = Pt(size)

        for n in card.key_numbers:
            add(f"• {n.label}: {n.value:+,.0f}  [{n.evidence.source}]", 16)
        add("Drivers: " + "; ".join(card.drivers), 14)
        add("Actions: " + "; ".join(card.actions), 14)
        so = ctx.get("signoff")
        add(f"Signed off by {so.approver}" if (so and so.approved) else "NOT signed off", 12)
        prs.save(out_path)
        return out_path


class OpenDesignRenderer:
    """Adapter for nexu-io/open-design (richer dashboards/PDF/PPTX/MP4, branded)."""
    def render(self, ctx: dict, out_path: str) -> str:
        raise NotImplementedError(
            "Open Design app not wired. Install nexu-io/open-design and point this adapter at "
            "it for branded dashboards/PPTX/PDF; until then renderer=builtin_html + PptxRenderer."
        )


def export_report(ctx: dict, out_dir: str = "out") -> dict:
    """Write the artifacts. Returns paths. Built-in HTML always; PPTX if python-pptx present."""
    Path(out_dir).mkdir(parents=True, exist_ok=True)
    html_path = str(Path(out_dir) / "card.html")
    Path(html_path).write_text(render_dashboard_html(ctx), encoding="utf-8")
    out = {"html": html_path, "pptx": None}
    pptx = PptxRenderer()
    if pptx.available():
        out["pptx"] = pptx.render(ctx, str(Path(out_dir) / "card.pptx"))
    return out
